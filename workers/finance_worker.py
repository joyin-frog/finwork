#!/usr/bin/env python3
import csv
import json
import os
import sys
from collections import defaultdict
from pathlib import Path

# ── CSV 分析域 ──────────────────────────────────────────────
def analyze_csv(path: Path):
    rows = []
    with path.open("r", encoding="utf-8-sig", newline="") as file:
        reader = csv.DictReader(file)
        fieldnames = reader.fieldnames or []
        for row in reader:
            rows.append(row)

    # 检测列名是否存在（列名不匹配时记入 column_warnings，不再静默返回 0）
    actual_cols = list(fieldnames)
    column_warnings: list[str] = []
    expected_cols = {
        "amount": "金额列（amount）",
        "category": "类目列（category）",
        "invoice_no": "发票号列（invoice_no）",
    }
    for col_key, col_label in expected_cols.items():
        if col_key not in actual_cols:
            known = "、".join(actual_cols) if actual_cols else "（空）"
            column_warnings.append(f"未找到{col_label}，已识别列：{known}")

    by_category_cents = defaultdict(int)
    warnings = []
    invoice_seen = set()
    for row in rows:
        amount = float(row.get("amount") or 0)
        category = row.get("category") or "未分类"
        invoice_no = (row.get("invoice_no") or "").strip()
        by_category_cents[category] += round(amount * 100)
        if amount <= 0:
            warnings.append({"invoice_no": invoice_no, "warning": "金额异常"})
        if invoice_no:
            if invoice_no in invoice_seen:
                warnings.append({"invoice_no": invoice_no, "warning": "发票号重复"})
            invoice_seen.add(invoice_no)

    return {
        "row_count": len(rows),
        "by_category": {k: v / 100 for k, v in by_category_cents.items()},
        "warnings": warnings,
        "column_warnings": column_warnings,
    }


# ── 文档解析域（xlsx/docx/pptx/pdf/OCR 提取） ──────────────
def _is_legacy_xls(path: Path) -> bool:
    """True for BIFF .xls. Never feed these to openpyxl (it only speaks OOXML)."""
    return path.suffix.lower() == ".xls"


def extract_xls(path: Path) -> str:
    """Read legacy .xls via xlrd — never openpyxl."""
    import xlrd

    book = xlrd.open_workbook(str(path))
    parts: list[str] = []
    for sheet in book.sheets():
        if sheet.nrows == 0:
            continue
        parts.append(f"## Sheet: {sheet.name}\n")
        max_cols = sheet.ncols
        for row_index in range(sheet.nrows):
            cells = []
            for col_index in range(max_cols):
                cell = sheet.cell(row_index, col_index)
                if cell.ctype == xlrd.XL_CELL_EMPTY:
                    cells.append("")
                else:
                    cells.append(str(cell.value))
            parts.append("| " + " | ".join(cells) + " |")
            if row_index == 0:
                parts.append("|" + "|".join(["---"] * max_cols) + "|")
        parts.append("")
    return "\n".join(parts)


def extract_xlsx(path: Path) -> str:
    if _is_legacy_xls(path):
        return extract_xls(path)

    import openpyxl

    wb = openpyxl.load_workbook(path, data_only=True)
    parts: list[str] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            continue
        parts.append(f"## Sheet: {sheet_name}\n")
        max_cols = max(len(row) for row in rows)
        for i, row in enumerate(rows):
            cells = [str(cell) if cell is not None else "" for cell in row]
            cells += [""] * (max_cols - len(cells))
            parts.append("| " + " | ".join(cells) + " |")
            if i == 0:
                parts.append("|" + "|".join(["---"] * max_cols) + "|")
        parts.append("")
    wb.close()
    return "\n".join(parts)


def inspect_xls(path: Path):
    """Inspect legacy .xls with xlrd (no formula engine; cached values only)."""
    import xlrd

    book = xlrd.open_workbook(str(path))
    workbook = {
        "file_name": path.name,
        "sheet_count": book.nsheets,
        "sheets": [],
        "format": "xls",
        "engine": "xlrd",
    }
    for sheet in book.sheets():
        header_row = [
            sheet.cell_value(0, col) if sheet.nrows > 0 else None
            for col in range(sheet.ncols)
        ]
        sample_rows = []
        for row_index in range(1, min(sheet.nrows, 8)):
            sample_rows.append([
                sheet.cell_value(row_index, col)
                for col in range(sheet.ncols)
            ])
        workbook["sheets"].append({
            "name": sheet.name,
            "rows": sheet.nrows,
            "columns": sheet.ncols,
            "headers": header_row,
            "sample_rows": sample_rows,
            "formula_count": 0,
            "formulas_sample": [],
            "merged_ranges": [],
            "frozen_panes": None,
            "auto_filter": None,
            "number_formats_sample": {},
        })
    return workbook


def inspect_excel(path: Path):
    if _is_legacy_xls(path):
        return inspect_xls(path)

    import openpyxl

    formula_wb = openpyxl.load_workbook(path, data_only=False, read_only=False)
    value_wb = openpyxl.load_workbook(path, data_only=True, read_only=True)

    workbook = {
        "file_name": path.name,
        "sheet_count": len(formula_wb.sheetnames),
        "sheets": [],
    }

    for sheet_name in formula_wb.sheetnames:
        ws = formula_wb[sheet_name]
        value_ws = value_wb[sheet_name]
        max_row = ws.max_row or 0
        max_column = ws.max_column or 0
        merged_ranges = [str(item) for item in ws.merged_cells.ranges]
        frozen_panes = str(ws.freeze_panes) if ws.freeze_panes else None

        header_row = []
        sample_rows = []
        formulas = []
        formula_count = 0
        formula_empty_cache_count = 0
        formula_errors = []
        layout_issues = []
        number_formats = {}

        for cell in next(ws.iter_rows(min_row=1, max_row=1, values_only=False), []):
            header_row.append(cell.value)

        sample_limit = min(max_row, 8)
        for row_index in range(2, sample_limit + 1):
            sample_rows.append([
                value_ws.cell(row=row_index, column=col_index).value
                for col_index in range(1, max_column + 1)
            ])

        merged_coordinates = {
            coordinate
            for merged_range in ws.merged_cells.ranges
            for row in ws[merged_range.coord]
            for cell in row
            for coordinate in [cell.coordinate]
        }
        for row in ws.iter_rows(values_only=False):
            for cell in row:
                if isinstance(cell.value, str) and cell.value.startswith("="):
                    formula_count += 1
                    cached_value = value_ws[cell.coordinate].value
                    if cached_value is None or cached_value == "":
                        formula_empty_cache_count += 1
                    if (
                        isinstance(cached_value, str)
                        and any(error in cached_value for error in (
                            "#DIV/0!", "#REF!", "#VALUE!", "#NAME?",
                            "#NULL!", "#NUM!", "#N/A",
                        ))
                        and len(formula_errors) < 100
                    ):
                        formula_errors.append({
                            "cell": cell.coordinate,
                            "cached_value": cached_value,
                        })
                    if len(formulas) < 25:
                        formulas.append({
                            "cell": cell.coordinate,
                            "formula": cell.value,
                            "cached_value": cached_value,
                        })
                if (
                    len(number_formats) < 40
                    and cell.number_format
                    and cell.number_format != "General"
                ):
                    number_formats[cell.coordinate] = cell.number_format
                if (
                    len(layout_issues) < 50
                    and isinstance(cell.value, str)
                    and not cell.value.startswith("=")
                    and len(cell.value.strip()) >= 16
                    and bool(cell.alignment.wrap_text)
                    and cell.coordinate not in merged_coordinates
                    and float(
                        ws.column_dimensions[cell.column_letter].width
                        or ws.sheet_format.defaultColWidth
                        or 8.43
                    ) <= 8
                ):
                    layout_issues.append({
                        "code": "narrow_wrapped_text",
                        "cell": cell.coordinate,
                        "text_length": len(cell.value.strip()),
                        "column_width": ws.column_dimensions[cell.column_letter].width or 8.43,
                    })

        workbook["sheets"].append({
            "name": sheet_name,
            "rows": max_row,
            "columns": max_column,
            "headers": header_row,
            "sample_rows": sample_rows,
            "formula_count": formula_count,
            "formula_empty_cache_count": formula_empty_cache_count,
            "formula_errors": formula_errors,
            "formulas_sample": formulas,
            "layout_issues": layout_issues,
            "merged_ranges": merged_ranges,
            "frozen_panes": frozen_panes,
            "auto_filter": str(ws.auto_filter.ref) if ws.auto_filter and ws.auto_filter.ref else None,
            "number_formats_sample": dict(list(number_formats.items())[:40]),
        })

    formula_wb.close()
    value_wb.close()
    return workbook


def extract_docx(path: Path) -> str:
    from docx import Document

    doc = Document(str(path))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs)


def extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        raise SystemExit("PPT 解析需要依赖未安装:pip install python-pptx")

    prs = Presentation(str(path))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        texts = [
            shape.text_frame.text
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        if texts:
            parts.append(f"--- Slide {i} ---\n" + "\n".join(texts))
    return "\n\n".join(parts)


def extract_pdf(path: Path) -> str:
    import pdfplumber

    parts: list[str] = []
    ocr_page_indices: set[int] = set()
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages):
            text = page.extract_text()
            if text and text.strip():
                parts.append(f"--- Page {i + 1} ---\n{text}")
            else:
                ocr_page_indices.add(i)
    # 文本层优先。OCR 只处理文本缺失的候选页，并设置硬上限，避免大 PDF
    # 被逐页 OCR 把 Agent/worker 卡死；调用方可通过环境变量提高上限。
    if ocr_page_indices:
        max_ocr_pages = max(0, int(os.environ.get("FINANCE_PDF_MAX_OCR_PAGES", "4")))
        candidates = sorted(ocr_page_indices)
        selected = set(candidates[:max_ocr_pages])
        ocr_text = _ocr_pdf_pages(path, selected) if selected else ""
        if ocr_text:
            parts.append(ocr_text)
        if len(candidates) > len(selected):
            parts.append(
                f"[PDF OCR truncated: {len(candidates) - len(selected)} pages omitted; "
                "use targeted page OCR after locating required fields.]"
            )
    return "\n\n".join(parts)


def _ocr_pdf_pages(path: Path, page_indices: set[int] | None = None) -> str:
    """图片型 PDF 兜底:pypdf 抽每页最大内嵌图(跳过 logo/印章小图)→ rapidocr OCR。"""
    from pypdf import PdfReader

    try:
        from rapidocr_onnxruntime import RapidOCR
    except ImportError:
        raise SystemExit("扫描件/图片型 PDF 的 OCR 需要依赖未安装:pip install rapidocr-onnxruntime")
    import numpy as np

    ocr = RapidOCR()
    reader = PdfReader(str(path))
    parts: list[str] = []
    for i, page in enumerate(reader.pages):
        if page_indices is not None and i not in page_indices:
            continue
        # 取该页面积最大的内嵌图 = 扫描主体,跳过 logo/印章/二维码等小图
        biggest = None
        biggest_area = 0
        for im in page.images:
            w, h = im.image.size
            if w * h > biggest_area:
                biggest_area = w * h
                biggest = im
        if biggest is None:
            continue
        arr = np.array(biggest.image.convert("RGB"))
        result, _ = ocr(arr, use_angle_cls=True)
        if not result:
            continue
        lines = sorted(result, key=lambda it: min(pt[1] for pt in it[0]))
        parts.append(f"--- Page {i + 1} (OCR) ---\n" + "\n".join(it[1] for it in lines))
    return "\n\n".join(parts)


def _ocr_image_file(path: Path) -> str:
    if not path.exists():
        raise FileNotFoundError(f"file not found: {path}")
    ext = path.suffix.lower()
    if ext not in (".png", ".jpg", ".jpeg", ".webp"):
        raise ValueError(f"unsupported image type: {ext}")

    try:
        from rapidocr_onnxruntime import RapidOCR
    except ImportError:
        raise RuntimeError("图片 OCR 需要依赖未安装:pip install rapidocr-onnxruntime")

    ocr = RapidOCR()
    # 手机拍的纸质单据常横拍/倒置;use_angle_cls 启用方向分类,自动摆正后再识别。
    result, _ = ocr(str(path), use_angle_cls=True)

    if not result:
        return ""

    # result 是 list of [box, text, score]; 按 box 左上角 y 坐标从上到下排序
    def _top_y(item):
        box = item[0]
        return min(pt[1] for pt in box)

    lines = sorted(result, key=_top_y)
    return "\n".join(item[1] for item in lines)


def cmd_ocr_image():
    if len(sys.argv) < 3:
        raise SystemExit("usage: finance_worker.py ocr-image <path>")
    try:
        print(_ocr_image_file(Path(sys.argv[2])))
    except Exception as error:
        raise SystemExit(str(error)) from error


def _extract_text_file(path: Path) -> str:
    if not path.exists():
        raise FileNotFoundError(f"file not found: {path}")
    ext = path.suffix.lower()
    if ext == ".xls":
        text = extract_xls(path)
    elif ext in (".xlsx", ".xlsm"):
        text = extract_xlsx(path)
    elif ext == ".docx":
        text = extract_docx(path)
    elif ext == ".pptx":
        text = extract_pptx(path)
    elif ext == ".pdf":
        text = extract_pdf(path)
    else:
        raise ValueError(f"unsupported file type: {ext}")
    return text


def cmd_extract_text():
    if len(sys.argv) < 3:
        raise SystemExit("usage: finance_worker.py extract-text <path>")
    try:
        print(_extract_text_file(Path(sys.argv[2])))
    except Exception as error:
        raise SystemExit(str(error)) from error


def cmd_document_server():
    """Long-lived NDJSON document server; one bounded request per line."""
    for raw in sys.stdin:
        request_id = None
        try:
            payload = json.loads(raw)
            request_id = payload.get("id")
            action = payload.get("action")
            file_path = payload.get("file_path")
            if not isinstance(file_path, str) or not file_path:
                raise ValueError("file_path must be a non-empty string")
            path = Path(file_path)
            if action == "extract-text":
                text = _extract_text_file(path)
            elif action == "ocr-image":
                text = _ocr_image_file(path)
            else:
                raise ValueError(f"unsupported action: {action}")
            result = {"ok": True, "text": text}
        except SystemExit as error:
            result = {"ok": False, "error": str(error)}
        except Exception as error:
            result = {"ok": False, "error": str(error)}
        result["id"] = request_id
        print(json.dumps(result, ensure_ascii=False), flush=True)


def cmd_inspect_excel():
    if len(sys.argv) < 3:
        raise SystemExit("usage: finance_worker.py inspect-excel <path>")
    path = Path(sys.argv[2])
    if not path.exists():
        raise SystemExit(f"file not found: {path}")
    if path.suffix.lower() not in (".xlsx", ".xlsm", ".xls"):
        raise SystemExit(f"unsupported file type: {path.suffix.lower()}")
    print(json.dumps(inspect_excel(path), ensure_ascii=False, indent=2, default=str))


def cmd_inspect_excel_cells():
    """inspect-excel-cells <path> <json-addresses> — read cached scalar values."""
    if len(sys.argv) < 4:
        raise SystemExit("usage: finance_worker.py inspect-excel-cells <path> <json-addresses>")
    path = Path(sys.argv[2])
    addresses = json.loads(sys.argv[3])
    import openpyxl

    wb = openpyxl.load_workbook(path, data_only=True, read_only=False)
    values = {}
    for qualified in addresses:
        if "!" not in qualified:
            values[qualified] = None
            continue
        sheet_name, address = qualified.rsplit("!", 1)
        values[qualified] = wb[sheet_name][address].value if sheet_name in wb.sheetnames else None
    wb.close()
    print(json.dumps({"ok": True, "values": values}, ensure_ascii=False, default=str))


def cmd_inspect_excel_formulas():
    """inspect-excel-formulas <path> <json-addresses> — read literal formulas."""
    if len(sys.argv) < 4:
        raise SystemExit("usage: finance_worker.py inspect-excel-formulas <path> <json-addresses>")
    path = Path(sys.argv[2])
    addresses = json.loads(sys.argv[3])
    import openpyxl

    wb = openpyxl.load_workbook(path, data_only=False, read_only=False)
    formulas = {}
    for qualified in addresses:
        if "!" not in qualified:
            formulas[qualified] = None
            continue
        sheet_name, address = qualified.rsplit("!", 1)
        if sheet_name not in wb.sheetnames:
            formulas[qualified] = None
            continue
        cell = wb[sheet_name][address]
        formulas[qualified] = cell.value if cell.data_type == "f" else None
    wb.close()
    print(json.dumps({"ok": True, "formulas": formulas}, ensure_ascii=False, default=str))


SPREADSHEET_NS = "http://schemas.openxmlformats.org/spreadsheetml/2006/main"
PACKAGE_REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
OFFICE_REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def _safe_xml_parser():
    """禁用实体解析:工作簿来自用户上传,不能给 XXE 留口子。"""
    from lxml import etree

    return etree.XMLParser(resolve_entities=False, no_network=True, huge_tree=False)


def _column_index(ref: str) -> int:
    """P7 → 16;用于把新建单元格插到正确的列序位置。"""
    index = 0
    for char in ref:
        if not char.isalpha():
            break
        index = index * 26 + (ord(char.upper()) - ord("A") + 1)
    return index


def _row_number(ref: str) -> int:
    digits = "".join(char for char in ref if char.isdigit())
    return int(digits) if digits else 0


def _sheet_xml_path(package: dict) -> dict:
    """
    工作表名 → zip 内 XML 路径。

    必须走 workbook.xml → r:id → workbook.xml.rels 这条链;按 sheet 顺序猜
    sheetN.xml 在删过工作表的文件上会错位,改错表比不改更危险。
    """
    from lxml import etree

    parser = _safe_xml_parser()
    workbook = etree.fromstring(package["xl/workbook.xml"], parser)
    rels = etree.fromstring(package["xl/_rels/workbook.xml.rels"], parser)
    target_by_id = {
        rel.get("Id"): rel.get("Target")
        for rel in rels.findall("{%s}Relationship" % PACKAGE_REL_NS)
    }
    mapping = {}
    for sheet in workbook.findall(".//{%s}sheets/{%s}sheet" % (SPREADSHEET_NS, SPREADSHEET_NS)):
        target = target_by_id.get(sheet.get("{%s}id" % OFFICE_REL_NS))
        if not target:
            continue
        normalized = target[1:] if target.startswith("/") else "xl/" + target.lstrip("./")
        mapping[sheet.get("name")] = normalized
    return mapping


def _find_or_create_cell(sheet_root, address: str):
    """定位单元格;行/列不存在时按序插入,保持 Excel 要求的升序。"""
    from lxml import etree

    sheet_data = sheet_root.find("{%s}sheetData" % SPREADSHEET_NS)
    if sheet_data is None:
        sheet_data = etree.SubElement(sheet_root, "{%s}sheetData" % SPREADSHEET_NS)
    row_number = _row_number(address)
    row = sheet_data.find("{%s}row[@r='%d']" % (SPREADSHEET_NS, row_number))
    if row is None:
        row = etree.Element("{%s}row" % SPREADSHEET_NS, r=str(row_number))
        following = [r for r in sheet_data if int(r.get("r", "0")) > row_number]
        if following:
            following[0].addprevious(row)
        else:
            sheet_data.append(row)
    cell = row.find("{%s}c[@r='%s']" % (SPREADSHEET_NS, address))
    if cell is None:
        cell = etree.Element("{%s}c" % SPREADSHEET_NS, r=address)
        column = _column_index(address)
        following = [c for c in row if _column_index(c.get("r", "")) > column]
        if following:
            following[0].addprevious(cell)
        else:
            row.append(cell)
    return cell


_NUMERIC_STRING_RE = None


def _numeric_string_pattern():
    """惰性编译:与其它正则辅助函数保持同样的按需 import 风格。"""
    global _NUMERIC_STRING_RE
    if _NUMERIC_STRING_RE is None:
        import re

        # 允许千分位逗号与小数点;含字母/百分号/货币符号一律不碰,交给下面的
        # inlineStr 分支——那些是明确的文本意图,不该被本函数猜成数字。
        _NUMERIC_STRING_RE = re.compile(r"^[+-]?(\d{1,3}(,\d{3})+|\d+)(\.\d+)?$")
    return _NUMERIC_STRING_RE


def _coerce_numeric_value(value):
    """
    JSON 字符串形式的纯数字 → 数值;其余原样返回,不确定就不碰。

    真实事故(2026-08-06,HISTORY-002):模型往 TB表 输入列写金额时序列化成了
    JSON 字符串("2850.8" 而不是 2850.8),被忠实写成 Excel 文本类型——SUM/ROUND
    等公式会静默跳过文本单元格,合计结果比实际少且不报任何错,是财务场景里
    最危险的一类静默错误(红线2:数值正确性)。

    取舍对齐 Excel 自身的默认行为:直接在单元格里敲数字,不预先设成文本格式,
    Excel 也会存成数值。**带前导零的字符串("007"、"0012")保持文本**——
    这类几乎总是账户代码/单号,数值化会丢掉前导零、改变精确匹配语义
    (TB表 A 列的科目代码就是这样故意存成文本的),这也正是 Excel 同样的
    默认行为(敲 007 会变成 7,除非预先设成文本格式)。
    """
    if not isinstance(value, str):
        return value
    text = value.strip()
    if not text or not _numeric_string_pattern().match(text):
        return value
    sign = ""
    body = text
    if body[0] in "+-":
        sign, body = body[0], body[1:]
    integer_part = body.split(".", 1)[0].replace(",", "")
    if len(integer_part) > 1 and integer_part[0] == "0":
        return value  # 前导零:多半是代码/单号,不当数值处理
    if len(integer_part) > 15:
        return value  # 超 15 位纯数字(统一社会信用代码/身份证号常见 18 位):
        # float 只能精确表示到约 2^53(16 位),再长会静默丢精度且改变原始
        # 文本语义,按文本处理更安全
    try:
        parsed = float(body.replace(",", ""))
    except ValueError:
        return value
    if sign == "-":
        parsed = -parsed
    if "." not in body and parsed == int(parsed):
        return int(parsed)
    return parsed


def _apply_cell_edit(cell, edit: dict) -> None:
    """按 edit 重写单元格内容;只动 <f>/<v>/<is> 与类型属性,样式 s= 原样保留。"""
    from lxml import etree

    for child in list(cell):
        cell.remove(child)
    cell.attrib.pop("t", None)
    if edit.get("clear"):
        return
    formula = edit.get("formula")
    if formula:
        node = etree.SubElement(cell, "{%s}f" % SPREADSHEET_NS)
        node.text = formula[1:] if formula.startswith("=") else formula
    value = edit.get("value")
    if value is None:
        return
    value = _coerce_numeric_value(value)
    if isinstance(value, bool):
        cell.set("t", "b")
        etree.SubElement(cell, "{%s}v" % SPREADSHEET_NS).text = "1" if value else "0"
    elif isinstance(value, (int, float)):
        etree.SubElement(cell, "{%s}v" % SPREADSHEET_NS).text = repr(value)
    elif formula:
        # 公式的字符串结果用 t="str",不进 sharedStrings。
        cell.set("t", "str")
        etree.SubElement(cell, "{%s}v" % SPREADSHEET_NS).text = str(value)
    else:
        # 内联字符串:避免改 sharedStrings.xml 牵动整册索引。
        cell.set("t", "inlineStr")
        is_node = etree.SubElement(cell, "{%s}is" % SPREADSHEET_NS)
        etree.SubElement(is_node, "{%s}t" % SPREADSHEET_NS).text = str(value)


def _drop_calc_chain(package: dict, names: list) -> None:
    """
    删掉 calcChain.xml 并摘除其引用。

    它是 Excel 的公式依赖缓存;改过单元格后留着旧链,Excel 会报「内容有问题」。
    删掉后 Excel 会自行重建。
    """
    from lxml import etree

    if "xl/calcChain.xml" not in package:
        return
    package.pop("xl/calcChain.xml")
    names.remove("xl/calcChain.xml")
    parser = _safe_xml_parser()
    rels_path = "xl/_rels/workbook.xml.rels"
    if rels_path in package:
        rels = etree.fromstring(package[rels_path], parser)
        for rel in rels.findall("{%s}Relationship" % PACKAGE_REL_NS):
            if (rel.get("Target") or "").endswith("calcChain.xml"):
                rels.remove(rel)
        package[rels_path] = etree.tostring(rels, xml_declaration=True, encoding="UTF-8", standalone=True)
    types_path = "[Content_Types].xml"
    if types_path in package:
        types = etree.fromstring(package[types_path], parser)
        for override in list(types):
            if (override.get("PartName") or "").endswith("calcChain.xml"):
                types.remove(override)
        package[types_path] = etree.tostring(types, xml_declaration=True, encoding="UTF-8", standalone=True)


MAX_EXCEL_ROW = 1_048_576
MAX_EXCEL_COL = 16_384
# 财务模板的依赖链很少超过几跳(TB表行 → 分类小计 → 报表科目 → 报表总计),
# 6 轮传播足够收敛,同时防止畸形循环引用把闭包算法拖入无限增长。
MAX_CLOSURE_PASSES = 6

_EXTERNAL_REF = None
_FULL_COL_REF = None
_CELL_OR_RANGE_REF = None


def _formula_ref_patterns():
    """惰性编译:避免模块顶层就 import re 影响其它命令的启动路径。"""
    global _EXTERNAL_REF, _FULL_COL_REF, _CELL_OR_RANGE_REF
    if _EXTERNAL_REF is not None:
        return _EXTERNAL_REF, _FULL_COL_REF, _CELL_OR_RANGE_REF
    import re

    sheet_prefix = r"(?:(?:'(?P<sq>[^']+)'|(?P<sb>[A-Za-z_一-鿿][\w.一-鿿]*))!)?"
    # 外部工作簿引用形如 [2]资产负债表!$C$7 或 '[2]资产负债表'!$C$7;先整体剔除,
    # 否则会被误判成本地表名"资产负债表"的引用——外部文件我们从不改,不该进闭包。
    _EXTERNAL_REF = re.compile(
        r"'?\[\d+\][^'!]*'?!\$?[A-Z]{1,3}\$?\d{1,7}(?::\$?[A-Z]{1,3}\$?\d{1,7})?"
    )
    # 整列引用,如 TB表!N:N —— SUMIFS/SUMIF 对科目余额表按整列取数是财务模板的常见写法,
    # 不支持这种形式会漏掉真实存在的依赖(2026-08-06 实测:C40 = SUMIFS(TB表!N:N,...))。
    # d1/d2 捕获列前的 $(有无决定共享公式展开时是否随列偏移)。
    _FULL_COL_REF = re.compile(
        sheet_prefix + r"(?P<d1>\$?)(?P<c1>[A-Z]{1,3})\$?:(?P<d2>\$?)(?P<c2>[A-Z]{1,3})(?!\d)"
    )
    _CELL_OR_RANGE_REF = re.compile(
        sheet_prefix
        + r"(?P<d1>\$?)(?P<c1>[A-Z]{1,3})(?P<d1r>\$?)(?P<r1>\d{1,7})"
        + r"(?::(?P<d2>\$?)(?P<c2>[A-Z]{1,3})(?P<d2r>\$?)(?P<r2>\d{1,7}))?"
    )
    return _EXTERNAL_REF, _FULL_COL_REF, _CELL_OR_RANGE_REF


def _extract_formula_refs(formula_text: str, own_sheet: str):
    """
    从公式文本里提取引用到的 (sheet, col_min, row_min, col_max, row_max) 列表。

    只做闭包判定需要的精度:不区分绝对/相对引用,不展开范围为具体单元格,
    命中一次假阳性(比如字符串里恰好长得像地址)最多让某个格子被多算一次
    重算候选——引擎给出的仍是该格子的真实计算结果,不会因此写错。
    """
    external_ref, full_col_ref, cell_or_range_ref = _formula_ref_patterns()
    cleaned = external_ref.sub(" ", formula_text)
    refs = []
    for match in full_col_ref.finditer(cleaned):
        sheet = match.group("sq") or match.group("sb") or own_sheet
        c1 = _column_index(match.group("c1"))
        c2 = _column_index(match.group("c2"))
        refs.append((sheet, min(c1, c2), 1, max(c1, c2), MAX_EXCEL_ROW))
    for match in cell_or_range_ref.finditer(cleaned):
        sheet = match.group("sq") or match.group("sb") or own_sheet
        c1 = _column_index(match.group("c1"))
        r1 = int(match.group("r1"))
        c2_raw, r2_raw = match.group("c2"), match.group("r2")
        c2 = _column_index(c2_raw) if c2_raw else c1
        r2 = int(r2_raw) if r2_raw else r1
        refs.append((sheet, min(c1, c2), min(r1, r2), max(c1, c2), max(r1, r2)))
    return refs


def _column_letters(index: int) -> str:
    """1 → 'A'，27 → 'AA'；`_column_index` 的逆运算。"""
    letters = ""
    index = max(1, index)
    while index > 0:
        index, remainder = divmod(index - 1, 26)
        letters = chr(65 + remainder) + letters
    return letters


def _offset_formula(formula_text: str, row_delta: int, col_delta: int) -> str:
    """
    展开共享公式(`<f t="shared">`):相对引用按 (行偏移,列偏移) 平移,`$` 锁定的维度不动。

    这是共享公式的核心语义——Excel 把结构相同的公式只存一份「主格」,其余「从格」
    只记一个索引,复制时相对地址自动跟着平移,加了 `$` 的部分保持不变。不展开的话,
    这些从格在 XML 里根本没有公式文本,会被当成普通值格漏掉(2026-08-06 实测:
    TB表 1735 个公式格里 443 个是这种从格,漏掉的往往正是财务模板里的行级汇总公式)。
    外部工作簿引用与 sheet 前缀原样保留,只有本地列/行地址参与平移。
    """
    if row_delta == 0 and col_delta == 0:
        return formula_text
    external_ref, full_col_ref, cell_or_range_ref = _formula_ref_patterns()

    externals: list[str] = []

    def stash_external(match: "re.Match[str]") -> str:
        externals.append(match.group(0))
        return "\x00%d\x00" % (len(externals) - 1)

    def sheet_prefix_of(match: "re.Match[str]") -> str:
        if match.group("sq"):
            return "'%s'!" % match.group("sq")
        if match.group("sb"):
            return "%s!" % match.group("sb")
        return ""

    def shift(token: str, dollar: str, delta: int, to_letters: bool) -> str:
        if dollar:
            return dollar + token
        if to_letters:
            return _column_letters(_column_index(token) + delta)
        return str(max(1, int(token) + delta))

    masked = external_ref.sub(stash_external, formula_text)

    def sub_full_col(match: "re.Match[str]") -> str:
        c1 = shift(match.group("c1"), match.group("d1"), col_delta, True)
        c2 = shift(match.group("c2"), match.group("d2"), col_delta, True)
        return "%s%s:%s" % (sheet_prefix_of(match), c1, c2)

    masked = full_col_ref.sub(sub_full_col, masked)

    def sub_cell(match: "re.Match[str]") -> str:
        prefix = sheet_prefix_of(match)
        c1 = shift(match.group("c1"), match.group("d1"), col_delta, True)
        r1 = shift(match.group("r1"), match.group("d1r"), row_delta, False)
        if match.group("c2"):
            c2 = shift(match.group("c2"), match.group("d2"), col_delta, True)
            r2 = shift(match.group("r2"), match.group("d2r"), row_delta, False)
            return "%s%s%s:%s%s" % (prefix, c1, r1, c2, r2)
        return "%s%s%s" % (prefix, c1, r1)

    masked = cell_or_range_ref.sub(sub_cell, masked)

    for index, external in enumerate(externals):
        masked = masked.replace("\x00%d\x00" % index, external)
    return masked


def _iter_formula_cells(root, sheet_name: str):
    """
    按 sheet 遍历 XML,产出 (地址, 公式文本含前导=, 引用列表)。

    两遍扫描:先收集共享公式(`t="shared"`)的「主格」(有完整文本、带 si 索引),
    再处理全部单元格——普通公式直接用自身文本;共享「从格」(同 si、无文本)按
    地址差对主格文本做 `_offset_formula` 展开,否则它们会被当成非公式格漏掉。
    """
    masters: dict[str, tuple[str, str]] = {}
    for cell in root.findall(".//{%s}c" % SPREADSHEET_NS):
        formula_node = cell.find("{%s}f" % SPREADSHEET_NS)
        if formula_node is None or formula_node.get("t") != "shared":
            continue
        si = formula_node.get("si")
        text = (formula_node.text or "").strip()
        address = cell.get("r")
        if si is not None and text and address:
            masters[si] = (address, text)

    for cell in root.findall(".//{%s}c" % SPREADSHEET_NS):
        formula_node = cell.find("{%s}f" % SPREADSHEET_NS)
        if formula_node is None:
            continue
        address = cell.get("r")
        if not address:
            continue
        raw = (formula_node.text or "").strip()
        if raw:
            text = "=" + raw
        elif formula_node.get("t") == "shared" and formula_node.get("si") in masters:
            master_address, master_text = masters[formula_node.get("si")]
            row_delta = _row_number(address) - _row_number(master_address)
            col_delta = _column_index(address) - _column_index(master_address)
            text = "=" + _offset_formula(master_text, row_delta, col_delta)
        else:
            continue
        yield address, text, _extract_formula_refs(text, sheet_name)


def _dependency_closure(all_formula_cells, changed_points, changed_keys):
    """
    从改动点出发,按引用关系传播,找出所有间接受影响的既有公式格。

    `changed_points` 是本轮直接写入的 (sheet, col, row);`changed_keys` 是同一批
    改动的 (sheet, address) 字符串键,用于把「被编辑的格子自己」排除出结果——
    它们已经在 `applied`/`formula_only` 里处理过,不该重复出现在下游清单里。
    返回 [(sheet, address, formula_text)],按发现顺序,即拓扑上由近到远。
    """
    frontier = set(changed_points)
    discovered_keys = set()
    stale = []
    for _ in range(MAX_CLOSURE_PASSES):
        if not frontier:
            break
        next_frontier = set()
        for sheet, address, text, refs in all_formula_cells:
            key = (sheet, address)
            if key in discovered_keys or key in changed_keys:
                continue
            hit = any(
                ref_sheet == point_sheet and c1 <= point_col <= c2 and r1 <= point_row <= r2
                for ref_sheet, c1, r1, c2, r2 in refs
                for point_sheet, point_col, point_row in frontier
            )
            if not hit:
                continue
            discovered_keys.add(key)
            stale.append((sheet, address, text))
            next_frontier.add((sheet, _column_index(address), _row_number(address)))
        frontier = next_frontier
    return stale


def _parse_solution_key(key: str):
    """
    `'[账簿.xlsx]财务报表'!P35` → `('财务报表', 'P35')`。

    formulas 的解按「书名+表名+地址」编码;只有解析出表名和地址才能对回我们
    关心的那几个格子。解析不出来的键(区域、命名区间、中间节点)一律跳过。
    """
    import re

    match = re.match(r"^'?\[[^\]]*\](?P<sheet>[^']*)'?!(?P<cell>\$?[A-Z]{1,3}\$?\d{1,7})$", key)
    if not match:
        return None
    return match.group("sheet"), match.group("cell").replace("$", "")


#  引擎在真实大模板上可能挂很久(见 _engine_cell_values 的超时熔断说明);
# 45s 远小于 agent 单次工具调用的合理耐心,超了就该放弃而不是拖死整个会话。
ENGINE_TIMEOUT_SECONDS = 45


def _engine_cell_values_sync(path, targets):
    """
    用 formulas 引擎求值,只取 targets 里的单元格。真正的计算逻辑,不含超时控制。

    刻意**不重算整册**的结果去覆盖文件:真实工作簿里大量公式指向别人机器上的
    外部链接(实测 2482 格里 1233 格算成错误值),整体覆盖会把原本有效的缓存
    值换成错误值。这里只回填目标格,且只在引擎真能算出来时才用。
    """
    try:
        import formulas
    except ImportError:
        return {}, "engine_not_installed", {}
    try:
        model = formulas.ExcelModel().loads(str(path)).finish()
        solution = model.calculate()
    except Exception as exc:  # 引擎对某些工作簿会直接抛错,不能连累交付
        return {}, "engine_failed: %s" % str(exc)[:200], {}

    lookup = {(sheet.lower(), cell) for sheet, cell in targets}
    values, skipped = {}, {}
    for key, node in solution.items():
        parsed = _parse_solution_key(str(key))
        if parsed is None:
            continue
        sheet, cell = parsed
        if (sheet.lower(), cell) not in lookup:
            continue
        try:
            raw = node.value[0, 0]
        except Exception:
            raw = node
        text = str(raw)
        if text.startswith("#"):
            skipped[(sheet.lower(), cell)] = "formula_error:%s" % text
            continue
        if hasattr(raw, "item"):
            raw = raw.item()  # numpy 标量 → Python 原生
        if isinstance(raw, str) and raw == "":
            # IFERROR(内层, "") 在内层够不到外部链接时就返回空串。
            # **不要**把它当计算结果写进缓存:在原作者那台有外部文件的机器上,
            # 这些格子是真实数字。写空串等于把「算不到」伪装成「算出来是空」。
            # 留空 → 断言标未验证 → 走人工校验,这是刻意的降级。
            skipped[(sheet.lower(), cell)] = "empty_result_likely_unreachable_source"
            continue
        values[(sheet.lower(), cell)] = raw
    return values, None, skipped


def _engine_cell_values_worker(path, targets, queue) -> None:
    """子进程入口:算完把结果丢进 queue,异常也转成消息而不是让子进程带着报错退出。"""
    try:
        values, note, skipped = _engine_cell_values_sync(path, targets)
        queue.put({"ok": True, "values": values, "note": note, "skipped": skipped})
    except Exception as exc:  # pragma: no cover - 兜底,不能让子进程静默消失
        queue.put({"ok": False, "error": "engine_crashed: %s" % str(exc)[:200]})


def _engine_cell_values(path, targets):
    """
    `_engine_cell_values_sync` 的超时熔断包装。

    2026-08-06 实测:真实合并报表模板(302 行 TB表 + 大量全列 SUMIFS 交叉引用)
    会让 `formulas` 库的图构建卡住——90 秒仍未跑完 `.finish()`,且没有任何
    进度反馈。这正是 HISTORY-002 两轮评测均在 40 分钟耗尽、finalize 从未
    成功的真实根因:此前的同步调用完全没有超时保护,一次挂起就吃光整个预算。

    用子进程执行 + 硬 join 超时:超时直接杀掉子进程,按「引擎不可用」处理——
    不阻塞调用方,原缓存保持不变,交由断言层标未验证(见 CONTEXT.md 不变量 2:
    能力不可用 ≠ 校验失败)。
    """
    import multiprocessing

    queue: "multiprocessing.Queue" = multiprocessing.Queue()
    process = multiprocessing.Process(
        target=_engine_cell_values_worker,
        args=(str(path), list(targets), queue),
    )
    process.start()
    process.join(ENGINE_TIMEOUT_SECONDS)
    if process.is_alive():
        process.terminate()
        process.join(5)
        if process.is_alive():
            process.kill()
            process.join(2)
        return {}, "engine_timeout(%ds)" % ENGINE_TIMEOUT_SECONDS, {}
    try:
        payload = queue.get_nowait()
    except Exception:
        return {}, "engine_crashed_no_result", {}
    if not payload.get("ok"):
        return {}, payload.get("error", "engine_failed"), {}
    return payload["values"], payload.get("note"), payload.get("skipped", {})


def _create_sheet(package: dict, names: list, sheet_name: str) -> str:
    """
    新建一张空工作表,返回它的 zip 内路径。

    要同时注册三处,少一处 Excel 就打不开:
      workbook.xml(表名+rId) / workbook.xml.rels(rId→路径) / [Content_Types].xml(part 类型)
    """
    from lxml import etree

    parser = _safe_xml_parser()
    workbook = etree.fromstring(package["xl/workbook.xml"], parser)
    rels = etree.fromstring(package["xl/_rels/workbook.xml.rels"], parser)
    types = etree.fromstring(package["[Content_Types].xml"], parser)

    used_ids = {rel.get("Id") for rel in rels.findall("{%s}Relationship" % PACKAGE_REL_NS)}
    index = 1
    while "rId%d" % index in used_ids:
        index += 1
    rel_id = "rId%d" % index
    seq = 1
    while "xl/worksheets/sheet%d.xml" % seq in package:
        seq += 1
    sheet_path = "xl/worksheets/sheet%d.xml" % seq

    sheet_ids = [int(s.get("sheetId", "0")) for s in workbook.findall(".//{%s}sheets/{%s}sheet" % (SPREADSHEET_NS, SPREADSHEET_NS))]
    sheets_node = workbook.find("{%s}sheets" % SPREADSHEET_NS)
    node = etree.SubElement(sheets_node, "{%s}sheet" % SPREADSHEET_NS)
    node.set("name", sheet_name)
    node.set("sheetId", str(max(sheet_ids, default=0) + 1))
    node.set("{%s}id" % OFFICE_REL_NS, rel_id)

    rel = etree.SubElement(rels, "{%s}Relationship" % PACKAGE_REL_NS)
    rel.set("Id", rel_id)
    rel.set("Type", "%s/worksheet" % OFFICE_REL_NS)
    rel.set("Target", "worksheets/sheet%d.xml" % seq)

    override = etree.SubElement(types, "{%s}Override" % types.nsmap[None])
    override.set("PartName", "/" + sheet_path)
    override.set(
        "ContentType",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.worksheet+xml",
    )

    package["xl/workbook.xml"] = etree.tostring(workbook, xml_declaration=True, encoding="UTF-8", standalone=True)
    package["xl/_rels/workbook.xml.rels"] = etree.tostring(rels, xml_declaration=True, encoding="UTF-8", standalone=True)
    package["[Content_Types].xml"] = etree.tostring(types, xml_declaration=True, encoding="UTF-8", standalone=True)
    package[sheet_path] = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n'
        '<worksheet xmlns="%s"><sheetData/></worksheet>' % SPREADSHEET_NS
    ).encode("utf-8")
    names.append(sheet_path)
    return sheet_path


def cmd_patch_workbook():
    """
    patch-workbook <src> <dst> <json-edits> — 在 XML 层就地改单元格。

    存在的理由:openpyxl 打开再保存会**丢掉全部公式缓存值**(实测 1164 → 0)。
    编辑用户已有的工作簿时那是不可接受的损失——外部链接的值再也算不回来。
    本命令只重写目标单元格,其余字节原样搬运。
    """
    import zipfile
    from lxml import etree

    if len(sys.argv) < 5:
        raise SystemExit("usage: finance_worker.py patch-workbook <src> <dst> <json-edits>")
    src, dst, edits = Path(sys.argv[2]), Path(sys.argv[3]), json.loads(sys.argv[4])

    with zipfile.ZipFile(src) as zin:
        names = zin.namelist()
        package = {name: zin.read(name) for name in names}

    sheet_paths = _sheet_xml_path(package)
    parser = _safe_xml_parser()
    roots, applied, missing, created_sheets = {}, [], [], []
    for edit in edits:
        sheet_name, address = edit.get("sheet"), edit.get("cell")
        path = sheet_paths.get(sheet_name)
        if (not path or path not in package) and edit.get("createSheet"):
            # 显式声明才建表:默认拒绝,避免把写错的表名静默变成一张新空表。
            path = _create_sheet(package, names, sheet_name)
            sheet_paths[sheet_name] = path
            created_sheets.append(sheet_name)
        if not path or path not in package:
            missing.append({"sheet": sheet_name, "cell": address, "reason": "sheet_not_found"})
            continue
        if path not in roots:
            roots[path] = etree.fromstring(package[path], parser)
        _apply_cell_edit(_find_or_create_cell(roots[path], address), edit)
        applied.append({"sheet": sheet_name, "cell": address})

    for path, root in roots.items():
        package[path] = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
    if applied:
        _drop_calc_chain(package, names)

    dst.parent.mkdir(parents=True, exist_ok=True)
    with zipfile.ZipFile(dst, "w", zipfile.ZIP_DEFLATED) as zout:
        for name in names:
            zout.writestr(name, package[name])

    # 只写了公式、没给结果的格子:读回是空的。模型写公式不带值是常态行为。
    formula_only = [
        edit for edit in edits
        if edit.get("formula") and edit.get("value") is None and not edit.get("clear")
        and any(a["sheet"] == edit.get("sheet") and a["cell"] == edit.get("cell") for a in applied)
    ]

    # 下游传播:被改单元格可能是既有公式的输入(填 TB表 一行,资产负债表的汇总公式
    # 就该跟着变)。只做「谁引用了改动点」的传播,不碰其余无关公式——
    # 2026-08-06 实测:HISTORY-002 填完 TB表 后,资产负债表汇总格(SUMIFS 整列引用、
    # 跨格 SUM)缓存仍是模板原始值,读不出真实数字,模型反复摸索直到超时。
    for sheet_name, path in sheet_paths.items():
        if path in package and path not in roots:
            roots[path] = etree.fromstring(package[path], parser)
    all_formula_cells = []
    for sheet_name, path in sheet_paths.items():
        if path not in roots:
            continue
        for address, text, refs in _iter_formula_cells(roots[path], sheet_name):
            all_formula_cells.append((sheet_name, address, text, refs))
    changed_keys = {(item["sheet"], item["cell"]) for item in applied}
    changed_points = {
        (item["sheet"], _column_index(item["cell"]), _row_number(item["cell"])) for item in applied
    }
    downstream = _dependency_closure(all_formula_cells, changed_points, changed_keys)

    # 合并两类待求值目标:模型新写的公式 + 因改动而失效的既有公式。
    # 用同一次引擎调用一起算,谁也不重复解析工作簿。
    recompute_formula, recompute_origin = {}, {}
    for edit in formula_only:
        key = (edit["sheet"], edit["cell"])
        formula_text = edit["formula"] if edit["formula"].startswith("=") else "=" + edit["formula"]
        recompute_formula[key] = formula_text
        recompute_origin[key] = "explicit"
    for sheet_name, address, text in downstream:
        key = (sheet_name, address)
        if key in recompute_formula:
            continue
        recompute_formula[key] = text
        recompute_origin[key] = "downstream"

    backfilled, engine_note = [], None
    unresolved_explicit, unresolved_downstream = [], []
    if recompute_formula:
        computed, engine_note, engine_skipped = _engine_cell_values(dst, set(recompute_formula.keys()))
        for (sheet_name, cell), formula_text in recompute_formula.items():
            # 引擎给不出有效值时:原缓存原样留着,不拿一个不确定的结果去覆盖它——
            # 「没算出来」和「算出来是错的」是两回事,前者才是这里唯一允许的降级。
            value = computed.get((sheet_name.lower(), cell)) if computed else None
            if value is None:
                target = unresolved_explicit if recompute_origin[(sheet_name, cell)] == "explicit" else unresolved_downstream
                target.append({"cell": "%s!%s" % (sheet_name, cell), "formula": formula_text})
                continue
            path = sheet_paths[sheet_name]
            if path not in roots:
                roots[path] = etree.fromstring(package[path], parser)
            cell_node = _find_or_create_cell(roots[path], cell)
            _apply_cell_edit(cell_node, {"formula": formula_text, "value": value})
            backfilled.append({
                "sheet": sheet_name,
                "cell": cell,
                "value": value,
                "reason": recompute_origin[(sheet_name, cell)],
            })
        if backfilled:
            for path in {sheet_paths[item["sheet"]] for item in backfilled}:
                package[path] = etree.tostring(
                    roots[path], xml_declaration=True, encoding="UTF-8", standalone=True
                )
            with zipfile.ZipFile(dst, "w", zipfile.ZIP_DEFLATED) as zout:
                for name in names:
                    zout.writestr(name, package[name])

    import openpyxl

    book = openpyxl.load_workbook(dst, data_only=False)
    cached_after = 0
    formula_total = 0
    values = openpyxl.load_workbook(dst, data_only=True)
    for sheet_name in book.sheetnames:
        for row in book[sheet_name].iter_rows():
            for cell in row:
                if not (isinstance(cell.value, str) and cell.value.startswith("=")):
                    continue
                formula_total += 1
                if values[sheet_name][cell.coordinate].value is not None:
                    cached_after += 1
    book.close()
    values.close()

    print(json.dumps({
        "ok": True,
        "applied": applied,
        "missing": missing,
        "createdSheets": created_sheets,
        "formulaCount": formula_total,
        "cachedValueCount": cached_after,
        "formulaOnlyCount": len(formula_only),
        # 新写入的公式 + 因改动而重算的既有公式,统一在这里;reason 区分来源。
        "backfilled": backfilled,
        "backfilledCount": len(backfilled),
        # 只写了公式、引擎也算不出来的:这些格子读回仍为空,需人工校验或装 LibreOffice。
        "unresolvedFormulaCells": [item["cell"] for item in unresolved_explicit],
        # 因改动而失效、但引擎也算不出来的既有公式:原缓存原样保留,标记待人工核对。
        "staleFormulas": unresolved_downstream[:200],
        "staleFormulaCount": len(unresolved_downstream),
        "engineNote": engine_note,
    }, ensure_ascii=False, default=str))


def cmd_compare_excel_allowlist():
    """compare-excel-allowlist <reference> <candidate> <sheet> <json-columns>."""
    if len(sys.argv) < 6:
        raise SystemExit(
            "usage: finance_worker.py compare-excel-allowlist "
            "<reference> <candidate> <sheet> <json-columns>"
        )
    reference_path = Path(sys.argv[2])
    candidate_path = Path(sys.argv[3])
    allowed_sheet = sys.argv[4]
    allowed_columns = {str(column).upper() for column in json.loads(sys.argv[5])}
    import openpyxl
    from openpyxl.utils import get_column_letter

    before_wb = openpyxl.load_workbook(reference_path, data_only=False, read_only=False)
    after_wb = openpyxl.load_workbook(candidate_path, data_only=False, read_only=False)
    before_sheets = before_wb.sheetnames
    after_sheets = after_wb.sheetnames
    disallowed = []
    changed_count = 0
    allowed_changed_count = 0
    if before_sheets != after_sheets:
        disallowed.append(
            "sheet-order:" + "|".join(before_sheets) + "=>" + "|".join(after_sheets)
        )

    def comparable(cell):
        if cell.data_type == "f":
            return ("formula", cell.value)
        value = cell.value
        if hasattr(value, "isoformat"):
            value = value.isoformat()
        return (cell.data_type, value)

    for sheet_name in dict.fromkeys([*before_sheets, *after_sheets]):
        if sheet_name not in before_wb.sheetnames or sheet_name not in after_wb.sheetnames:
            continue
        before = before_wb[sheet_name]
        after = after_wb[sheet_name]
        max_row = max(before.max_row or 0, after.max_row or 0)
        max_column = max(before.max_column or 0, after.max_column or 0)
        for row in range(1, max_row + 1):
            for column in range(1, max_column + 1):
                if comparable(before.cell(row, column)) == comparable(after.cell(row, column)):
                    continue
                changed_count += 1
                address = f"{get_column_letter(column)}{row}"
                if sheet_name == allowed_sheet and get_column_letter(column) in allowed_columns:
                    allowed_changed_count += 1
                elif len(disallowed) < 100:
                    disallowed.append(f"{sheet_name}!{address}")
    before_wb.close()
    after_wb.close()
    print(json.dumps({
        "ok": True,
        "changedCount": changed_count,
        "allowedChangedCount": allowed_changed_count,
        "disallowedChanges": disallowed,
    }, ensure_ascii=False))


def cmd_convert_xls():
    """convert-xls <input.xls> <output.xlsx> — xlrd read → openpyxl write."""
    if len(sys.argv) < 4:
        raise SystemExit("usage: finance_worker.py convert-xls <input.xls> <output.xlsx>")
    src = Path(sys.argv[2])
    dst = Path(sys.argv[3])
    if not src.exists():
        print(json.dumps({"ok": False, "error": f"file not found: {src}"}, ensure_ascii=False))
        return
    if src.suffix.lower() != ".xls":
        print(json.dumps({"ok": False, "error": "input must be .xls"}, ensure_ascii=False))
        return
    try:
        import xlrd
        import openpyxl
    except ImportError as e:
        print(json.dumps({"ok": False, "error": f"missing dependency: {e}"}, ensure_ascii=False))
        return

    book = xlrd.open_workbook(str(src))
    wb = openpyxl.Workbook()
    default = wb.active
    wb.remove(default)
    for sheet in book.sheets():
        ws = wb.create_sheet(title=sheet.name[:31] or "Sheet")
        for row_index in range(sheet.nrows):
            for col_index in range(sheet.ncols):
                cell = sheet.cell(row_index, col_index)
                value = cell.value
                if cell.ctype == xlrd.XL_CELL_DATE:
                    try:
                        value = xlrd.xldate_as_datetime(cell.value, book.datemode)
                    except Exception:
                        pass
                ws.cell(row=row_index + 1, column=col_index + 1, value=value)
    dst.parent.mkdir(parents=True, exist_ok=True)
    wb.save(str(dst))
    print(json.dumps({"ok": True, "outputPath": str(dst)}, ensure_ascii=False))


def _fixture_dir() -> Path:
    """Spreadsheet probe fixtures (repo tests/fixtures/spreadsheet)."""
    env = os.environ.get("FINANCE_AGENT_SPREADSHEET_FIXTURES")
    if env:
        return Path(env)
    return Path(__file__).resolve().parent.parent / "tests" / "fixtures" / "spreadsheet"


def cmd_probe_spreadsheet():
    """Behavioral package probe: xlrd .xls read + openpyxl xlsx roundtrip."""
    import tempfile

    result = {
        "read": {"csv": False, "xlsx": False, "xlsm": False, "xls": False},
        "write": {"xlsx": False},
        "ok": False,
    }
    fixtures = _fixture_dir()
    try:
        import xlrd
        xls_path = fixtures / "legacy-input.xls"
        if xls_path.exists():
            book = xlrd.open_workbook(str(xls_path))
            sheet = book.sheet_by_index(0)
            # Fixture contract: A1 == "Name", B2 == 42
            a1 = sheet.cell_value(0, 0)
            b2 = sheet.cell_value(1, 1) if sheet.nrows > 1 and sheet.ncols > 1 else None
            result["read"]["xls"] = a1 == "Name" and float(b2) == 42.0
    except Exception as e:
        result["xls_error"] = str(e)

    try:
        import openpyxl
        with tempfile.TemporaryDirectory() as td:
            out = Path(td) / "probe.xlsx"
            wb = openpyxl.Workbook()
            ws = wb.active
            ws["A1"] = "hello"
            ws["A2"] = 1
            ws["A3"] = 2
            ws["A4"] = "=SUM(A2:A3)"
            wb.save(out)
            result["write"]["xlsx"] = out.exists() and out.stat().st_size > 0
            rb = openpyxl.load_workbook(out, data_only=False)
            result["read"]["xlsx"] = rb.active["A1"].value == "hello"
            rb.close()
            xlsm = fixtures / "macro-input.xlsm"
            if xlsm.exists():
                try:
                    mb = openpyxl.load_workbook(xlsm, keep_vba=True, data_only=True)
                    result["read"]["xlsm"] = len(mb.sheetnames) >= 1
                    mb.close()
                except Exception:
                    result["read"]["xlsm"] = False
            else:
                result["read"]["xlsm"] = result["read"]["xlsx"]
        import pandas as pd
        with tempfile.TemporaryDirectory() as td:
            csv_path = Path(td) / "t.csv"
            csv_path.write_text("a,b\n1,2\n", encoding="utf-8")
            df = pd.read_csv(csv_path)
            result["read"]["csv"] = len(df) == 1 and int(df.iloc[0]["a"]) == 1
    except Exception as e:
        result["xlsx_error"] = str(e)

    result["ok"] = (
        result["read"]["xls"]
        and result["read"]["xlsx"]
        and result["write"]["xlsx"]
        and result["read"]["csv"]
    )
    print(json.dumps(result, ensure_ascii=False))


def cmd_probe_recalc():
    """If LibreOffice executable given, recalc =SUM(A1:A2) and expect 3."""
    if len(sys.argv) < 3:
        print(json.dumps({"ok": False, "error": "usage: probe-recalc <soffice>"}, ensure_ascii=False))
        return
    soffice = sys.argv[2]
    import tempfile
    import subprocess
    import openpyxl

    with tempfile.TemporaryDirectory() as td:
        td_path = Path(td)
        profile = td_path / "lo-profile"
        profile.mkdir()
        xlsx = td_path / "sum.xlsx"
        wb = openpyxl.Workbook()
        ws = wb.active
        ws["A1"] = 1
        ws["A2"] = 2
        ws["A3"] = "=SUM(A1:A2)"
        wb.save(xlsx)
        wb.close()
        out_dir = td_path / "out"
        out_dir.mkdir()
        try:
            proc = subprocess.run(
                [
                    soffice,
                    "--headless",
                    "--norestore",
                    "--nolockcheck",
                    f"-env:UserInstallation=file://{profile}",
                    "--convert-to",
                    "xlsx",
                    "--outdir",
                    str(out_dir),
                    str(xlsx),
                ],
                capture_output=True,
                text=True,
                timeout=90,
                env={**os.environ, "SAL_USE_VCLPLUGIN": "svp"},
            )
        except subprocess.TimeoutExpired:
            print(json.dumps({"ok": False, "error": "recalc_timeout"}, ensure_ascii=False))
            return
        if proc.returncode != 0:
            print(json.dumps({
                "ok": False,
                "error": proc.stderr or proc.stdout or f"exit {proc.returncode}",
            }, ensure_ascii=False))
            return
        produced = list(out_dir.glob("*.xlsx"))
        target = produced[0] if produced else xlsx
        rb = openpyxl.load_workbook(target, data_only=True)
        value = rb.active["A3"].value
        rb.close()
        print(json.dumps({"ok": value == 3, "value": value}, ensure_ascii=False))


def cmd_recalc_xlsx():
    """recalc-xlsx <xlsx> <soffice> [timeout_seconds] — working-copy friendly helper."""
    if len(sys.argv) < 4:
        print(json.dumps({"ok": False, "error": "usage: recalc-xlsx <xlsx> <soffice> [timeout]"}, ensure_ascii=False))
        return
    xlsx = Path(sys.argv[2])
    soffice = sys.argv[3]
    timeout = int(sys.argv[4]) if len(sys.argv) > 4 else 60
    import tempfile
    import subprocess
    import openpyxl
    import shutil

    if not xlsx.exists():
        print(json.dumps({"ok": False, "error": f"file not found: {xlsx}"}, ensure_ascii=False))
        return

    with tempfile.TemporaryDirectory() as td:
        td_path = Path(td)
        profile = td_path / "lo-profile"
        profile.mkdir()
        work = td_path / xlsx.name
        shutil.copy2(xlsx, work)
        out_dir = td_path / "out"
        out_dir.mkdir()
        try:
            proc = subprocess.run(
                [
                    soffice,
                    "--headless",
                    "--norestore",
                    "--nolockcheck",
                    f"-env:UserInstallation=file://{profile}",
                    "--convert-to",
                    "xlsx",
                    "--outdir",
                    str(out_dir),
                    str(work),
                ],
                capture_output=True,
                text=True,
                timeout=timeout,
                env={**os.environ, "SAL_USE_VCLPLUGIN": "svp"},
            )
        except subprocess.TimeoutExpired:
            print(json.dumps({"ok": False, "error": "recalc_timeout"}, ensure_ascii=False))
            return
        if proc.returncode != 0:
            print(json.dumps({
                "ok": False,
                "error": proc.stderr or proc.stdout or f"exit {proc.returncode}",
            }, ensure_ascii=False))
            return
        produced = list(out_dir.glob("*.xlsx"))
        if not produced:
            print(json.dumps({"ok": False, "error": "no output workbook"}, ensure_ascii=False))
            return
        shutil.copy2(produced[0], xlsx)
        formula_wb = openpyxl.load_workbook(xlsx, data_only=False)
        formula_count = 0
        for sheet_name in formula_wb.sheetnames:
            ws = formula_wb[sheet_name]
            for row in ws.iter_rows(values_only=True):
                for value in row:
                    if isinstance(value, str) and value.startswith("="):
                        formula_count += 1
        formula_wb.close()
        print(json.dumps({"ok": True, "formulaCount": formula_count}, ensure_ascii=False))


# ── 固定导出域（版本化防覆盖） ─────────────────────────────
def _next_versioned_path(path: Path) -> Path:
    """path 已存在则返回同目录下 stem_v2/_v3… 的首个空位,用于「不覆盖上一版产物」。"""
    if not path.exists():
        return path
    n = 2
    while True:
        cand = path.with_name(f"{path.stem}_v{n}{path.suffix}")
        if not cand.exists():
            return cand
        n += 1


# ── demo 数据与环境自检域 ───────────────────────────────────
def demo():
    demo_file = get_demo_data_path()
    demo_file.parent.mkdir(parents=True, exist_ok=True)
    if not demo_file.exists():
        demo_file.write_text(
            "employee,expense_date,invoice_no,category,amount\n"
            "张敏,2026-05-12,INV-001,交通,380\n"
            "李哲,2026-05-14,INV-002,招待,1680\n"
            "王岚,2026-05-15,INV-001,办公,260\n",
            encoding="utf-8",
        )
    return analyze_csv(demo_file)


def get_demo_data_path():
    if os.environ.get("FINANCE_AGENT_DEMO_DATA_PATH"):
        return Path(os.environ["FINANCE_AGENT_DEMO_DATA_PATH"])
    if os.environ.get("FINANCE_AGENT_APP_DATA_DIR"):
        return Path(os.environ["FINANCE_AGENT_APP_DATA_DIR"]) / "demo_reimbursements.csv"
    if sys.platform == "darwin":
        return Path.home() / "Library" / "Application Support" / "finance-agent" / "demo_reimbursements.csv"
    if sys.platform == "win32":
        return Path(os.environ.get("APPDATA", Path.home() / "AppData" / "Roaming")) / "finance-agent" / "demo_reimbursements.csv"
    return Path(os.environ.get("XDG_DATA_HOME", Path.home() / ".local" / "share")) / "finance-agent" / "demo_reimbursements.csv"


def cmd_selfcheck():
    # 首启环境自检:报告 Python 版本与关键依赖是否就位(供桌面端 doctor 用人话提示用户)。
    deps = [
        "openpyxl",
        "pandas",
        "pdfplumber",
        "xlsxwriter",
        "xlrd",
        "pypdf",
        "reportlab",
        "docx",
        "pptx",
        "lxml",
        "PIL",
        "defusedxml",
        "pdf2image",
        "markitdown",
    ]
    found = {}
    missing = []
    for name in deps:
        try:
            module = __import__(name)
            found[name] = getattr(module, "__version__", None) or getattr(module, "__VERSION__", "unknown")
        except Exception:
            missing.append(name)
    print(json.dumps({
        "python": sys.version.split()[0],
        "deps": found,
        "missing": missing,
        "ok": len(missing) == 0,
    }, ensure_ascii=False))


def _force_utf8_stdio():
    """把 stdin/stdout/stderr 统一改成 UTF-8。Windows 中文系统默认 cp936/GBK:Node 按 UTF-8 把
    JSON 写进 stdin → sys.stdin.read() 误解码出游离代理(\\udcXX)→ print 报 "surrogates not
    allowed"。在读 stdin / 任何 print 之前调用，保证固定 worker
    被直接运行/测试(没设那个 env)也正确。"""
    for stream in (sys.stdin, sys.stdout, sys.stderr):
        try:
            stream.reconfigure(encoding="utf-8")  # type: ignore[union-attr]
        except (AttributeError, ValueError):
            pass  # 非常规流(如已被替换为 StringIO)忽略


# ── xlsx 导出域（凭证 / 工资条） ────────────────────────────
def cmd_export_voucher_xlsx():
    """export-voucher-xlsx: 从 stdin 读 JSON payload,生成三 sheet 金蝶对照清单 xlsx。
    JSON 结构:
      { outputPath: str, vouchers: [...], skipped?: [...] }
    输出到 stdout: JSON { filePath: str }
    """
    import openpyxl

    payload = json.loads(sys.stdin.read())
    # 防覆盖:同名已存在则版本化为 _v2/_v3…,
    # 否则第二次导出会原地改写上一份交付物(附件按路径去重,旧回答的附件被静默篡改)
    output_path = str(_next_versioned_path(Path(payload["outputPath"])))
    vouchers = payload.get("vouchers", [])
    skipped = payload.get("skipped", [])
    # needs_confirm 只进 sheet3,不混入对照清单与汇总口径
    confirmed = [v for v in vouchers if v.get("status") != "needs_confirm"]
    needs_confirm = [v for v in vouchers if v.get("status") == "needs_confirm"]

    wb = openpyxl.Workbook()

    # ── Sheet 1: 对照清单(auto/confirmed 凭证) ──
    ws1 = wb.active
    ws1.title = "对照清单"
    header1 = ["日期", "凭证字", "摘要", "科目编码", "科目全名", "核算维度类型", "核算维度值", "借方金额", "贷方金额"]
    ws1.append(header1)
    for v in confirmed:
        date = v.get("date", "")
        voucher_word = v.get("voucherWord", "记")
        for line in v.get("lines", []):
            ws1.append([
                date,
                voucher_word,
                line.get("summary", ""),
                line.get("account", ""),
                line.get("accountName", ""),
                line.get("dimensionType", ""),
                line.get("dimensionValue", ""),
                line.get("debitYuan", None),
                line.get("creditYuan", None),
            ])

    # ── Sheet 2: 汇总(按文件、按科目笔数金额小计) ──
    ws2 = wb.create_sheet("汇总")
    ws2.append(["维度", "科目编码", "科目全名", "借方合计", "贷方合计", "笔数"])
    # 按科目汇总
    from collections import defaultdict
    by_account: dict = defaultdict(lambda: {"name": "", "debit": 0.0, "credit": 0.0, "count": 0})
    for v in confirmed:
        for line in v.get("lines", []):
            code = line.get("account", "")
            by_account[code]["name"] = line.get("accountName", code)
            by_account[code]["debit"] += line.get("debitYuan") or 0
            by_account[code]["credit"] += line.get("creditYuan") or 0
            by_account[code]["count"] += 1
    for code, info in sorted(by_account.items()):
        ws2.append(["科目", code, info["name"], round(info["debit"], 2), round(info["credit"], 2), info["count"]])
    # 按文件汇总
    ws2.append([])
    ws2.append(["文件", "文件名", "", "借方合计", "贷方合计", ""])
    for v in confirmed:
        file_debit = sum(line.get("debitYuan") or 0 for line in v.get("lines", []))
        file_credit = sum(line.get("creditYuan") or 0 for line in v.get("lines", []))
        ws2.append(["文件", v.get("file", ""), "", round(file_debit, 2), round(file_credit, 2), ""])

    # ── Sheet 3: 待确认与跳过 ──
    ws3 = wb.create_sheet("待确认与跳过")
    if not needs_confirm and not skipped:
        ws3.append(["无"])
    else:
        ws3.append(["类型", "文件", "摘要", "金额(元)", "原因"])
        for v in needs_confirm:
            issues = "; ".join(v.get("issues", []))
            total = sum((line.get("debitYuan") or 0) for line in v.get("lines", []))
            ws3.append(["待确认", v.get("file", ""), "", round(total, 2), issues])
        for s in skipped:
            ws3.append(["跳过", s.get("file", ""), s.get("summary", ""), s.get("amountYuan", ""), s.get("reason", "")])

    wb.save(output_path)
    print(json.dumps({"filePath": output_path}, ensure_ascii=False))


def cmd_export_payslips_xlsx():
    """export-payslips-xlsx: 从 stdin 读 JSON payload，生成工资明细 xlsx。
    JSON 结构:
      { outputPath: str, year: int, month: int,
        rows: [{employeeName, grossPay, socialInsurance, housingFund,
                specialDeduction, taxCurrent, netPay}, ...],
        totals: {grossPay, socialInsurance, housingFund, specialDeduction,
                 taxCurrent, netPay} }
    输出到 stdout: JSON { filePath: str }
    """
    import openpyxl

    payload = json.loads(sys.stdin.read())
    # 防覆盖：同名已存在则版本化为 _v2/_v3…（照抄 cmd_export_voucher_xlsx，
    # 显式调 _next_versioned_path，不依赖 _guarded_save 猴子补丁）
    output_path = str(_next_versioned_path(Path(payload["outputPath"])))
    year = payload.get("year", "")
    month = payload.get("month", "")
    rows = payload.get("rows", [])
    totals = payload.get("totals", {})

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "工资明细"

    # 第 1 行：标注期间与状态
    ws.append([f"{year}年{month}月工资明细 · 已确认"])

    # 表头
    header = ["姓名", "税前工资", "五险(个人)", "公积金(个人)", "专项附加扣除", "本期个税", "实发工资"]
    ws.append(header)

    # 每人一行
    for r in rows:
        ws.append([
            r.get("employeeName", ""),
            r.get("grossPay", 0),
            r.get("socialInsurance", 0),
            r.get("housingFund", 0),
            r.get("specialDeduction", 0),
            r.get("taxCurrent", 0),
            r.get("netPay", 0),
        ])

    # 合计行
    ws.append([
        "合计",
        totals.get("grossPay", 0),
        totals.get("socialInsurance", 0),
        totals.get("housingFund", 0),
        totals.get("specialDeduction", 0),
        totals.get("taxCurrent", 0),
        totals.get("netPay", 0),
    ])

    wb.save(output_path)
    print(json.dumps({"filePath": output_path}, ensure_ascii=False))


# ── embedding 域（语义检索，WP12） ──────────────────────────
_embedding_runtime_cache = {}


def _load_embedding_runtime(model_dir):
    import os as _os
    onnx_path = _os.path.join(model_dir, "model_quantized.onnx")
    tokenizer_path = _os.path.join(model_dir, "tokenizer.json")
    if not _os.path.exists(onnx_path) or not _os.path.exists(tokenizer_path):
        raise RuntimeError("model_not_found")
    cached = _embedding_runtime_cache.get(model_dir)
    if cached is not None:
        return cached
    try:
        from tokenizers import Tokenizer  # type: ignore
        import onnxruntime as ort  # type: ignore
        import numpy as np  # type: ignore
    except ImportError as e:
        raise RuntimeError(f"import_error: {e}") from e
    tokenizer = Tokenizer.from_file(tokenizer_path)
    tokenizer.enable_padding(pad_id=0, pad_token="[PAD]", length=512)
    tokenizer.enable_truncation(max_length=512)
    session = ort.InferenceSession(onnx_path, providers=["CPUExecutionProvider"])
    runtime = (tokenizer, session, np)
    _embedding_runtime_cache[model_dir] = runtime
    return runtime


def _embed_payload(payload):
    texts = payload.get("texts", [])
    model_dir = payload.get("model_dir", "")
    if not isinstance(texts, list) or any(not isinstance(text, str) for text in texts):
        return {"ok": False, "error": "invalid_texts"}
    if not texts:
        return {"ok": True, "dim": 512, "vectors": []}
    try:
        tokenizer, session, np = _load_embedding_runtime(model_dir)
        encodings = tokenizer.encode_batch(texts)
        input_ids = np.array([encoding.ids for encoding in encodings], dtype=np.int64)
        attention_mask = np.array([encoding.attention_mask for encoding in encodings], dtype=np.int64)
        token_type_ids = np.zeros_like(input_ids, dtype=np.int64)
        outputs = session.run(None, {
            "input_ids": input_ids,
            "attention_mask": attention_mask,
            "token_type_ids": token_type_ids,
        })
        token_embeddings = outputs[0]
        mask_expanded = attention_mask[:, :, None].astype(np.float32)
        embeddings = (token_embeddings * mask_expanded).sum(axis=1) / mask_expanded.sum(axis=1).clip(min=1e-9)
        embeddings = embeddings / np.linalg.norm(embeddings, axis=1, keepdims=True).clip(min=1e-9)
        vectors = embeddings.tolist()
        return {"ok": True, "dim": len(vectors[0]) if vectors else 512, "vectors": vectors}
    except Exception as e:
        message = str(e)
        if message == "model_not_found" or message.startswith("import_error:"):
            return {"ok": False, "error": message}
        return {"ok": False, "error": f"embed_error: {message}"}


def cmd_embed_texts():
    """Single-request compatibility command. The persistent pool uses embed-server."""
    try:
        payload = json.loads(sys.stdin.read())
    except Exception as e:
        print(json.dumps({"ok": False, "error": f"invalid_json: {e}"}, ensure_ascii=False))
        return
    print(json.dumps(_embed_payload(payload), ensure_ascii=False))


def cmd_embed_server():
    """Long-lived NDJSON embedding server; one request and one response per line."""
    for raw in sys.stdin:
        request_id = None
        try:
            payload = json.loads(raw)
            request_id = payload.get("id")
            result = _embed_payload(payload)
        except Exception as e:
            result = {"ok": False, "error": f"invalid_json: {e}"}
        result["id"] = request_id
        print(json.dumps(result, ensure_ascii=False), flush=True)


# ── 命令分发入口 ────────────────────────────────────────────
def main():
    _force_utf8_stdio()
    if len(sys.argv) >= 2 and sys.argv[1] == "--selfcheck":
        cmd_selfcheck()
        return
    if len(sys.argv) >= 2 and sys.argv[1] == "demo":
        print(json.dumps(demo(), ensure_ascii=False, indent=2))
        return
    if len(sys.argv) >= 2 and sys.argv[1] == "ocr-image":
        cmd_ocr_image()
        return
    if len(sys.argv) >= 2 and sys.argv[1] == "extract-text":
        cmd_extract_text()
        return
    if len(sys.argv) >= 2 and sys.argv[1] == "document-server":
        cmd_document_server()
        return
    if len(sys.argv) >= 2 and sys.argv[1] == "inspect-excel":
        cmd_inspect_excel()
        return
    if len(sys.argv) >= 2 and sys.argv[1] == "inspect-excel-cells":
        cmd_inspect_excel_cells()
        return
    if len(sys.argv) >= 2 and sys.argv[1] == "inspect-excel-formulas":
        cmd_inspect_excel_formulas()
        return
    if len(sys.argv) >= 2 and sys.argv[1] == "compare-excel-allowlist":
        cmd_compare_excel_allowlist()
        return
    if len(sys.argv) >= 2 and sys.argv[1] == "patch-workbook":
        cmd_patch_workbook()
        return
    if len(sys.argv) >= 2 and sys.argv[1] == "convert-xls":
        cmd_convert_xls()
        return
    if len(sys.argv) >= 2 and sys.argv[1] == "probe-spreadsheet":
        cmd_probe_spreadsheet()
        return
    if len(sys.argv) >= 2 and sys.argv[1] == "probe-recalc":
        cmd_probe_recalc()
        return
    if len(sys.argv) >= 2 and sys.argv[1] == "recalc-xlsx":
        cmd_recalc_xlsx()
        return
    if len(sys.argv) == 3 and sys.argv[1] == "analyze-csv":
        print(json.dumps(analyze_csv(Path(sys.argv[2])), ensure_ascii=False, indent=2))
        return
    if len(sys.argv) >= 2 and sys.argv[1] == "export-voucher-xlsx":
        cmd_export_voucher_xlsx()
        return
    if len(sys.argv) >= 2 and sys.argv[1] == "export-payslips-xlsx":
        cmd_export_payslips_xlsx()
        return
    if len(sys.argv) >= 2 and sys.argv[1] == "embed-texts":
        cmd_embed_texts()
        return
    if len(sys.argv) >= 2 and sys.argv[1] == "embed-server":
        cmd_embed_server()
        return
    raise SystemExit(
        "usage: finance_worker.py --selfcheck | demo | analyze-csv <path> | extract-text <path> | document-server | inspect-excel <path> | inspect-excel-cells <path> <json-addresses> | inspect-excel-formulas <path> <json-addresses> | compare-excel-allowlist <reference> <candidate> <sheet> <json-columns> | patch-workbook <src> <dst> <json-edits> | convert-xls <xls> <xlsx> | probe-spreadsheet | probe-recalc <soffice> | recalc-xlsx <xlsx> <soffice> [timeout] | ocr-image <path> | export-voucher-xlsx | export-payslips-xlsx | embed-texts | embed-server"
    )


if __name__ == "__main__":
    main()
